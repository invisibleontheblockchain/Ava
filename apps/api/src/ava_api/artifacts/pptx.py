"""Phase 1 deliverable: basic .pptx from Note block content (Document 2 Table 11)."""

from __future__ import annotations

import uuid
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def generate_pptx_from_text(*, title: str, body: str, output_dir: Path) -> uuid.UUID:
    output_dir.mkdir(parents=True, exist_ok=True)
    artifact_id = uuid.uuid4()
    path = output_dir / f"{artifact_id}.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = title[:200]
    tf = slide.placeholders[1].text_frame
    tf.text = body[:8000]
    for para in body[8000:16000].split("\n\n")[:5]:
        if not para.strip():
            continue
        p = tf.add_paragraph()
        p.text = para.strip()[:2000]
        p.font.size = Pt(14)

    buffer = BytesIO()
    prs.save(buffer)
    path.write_bytes(buffer.getvalue())
    return artifact_id
